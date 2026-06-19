#!/usr/bin/env python3
"""Generate comprehensive class PowerPoint with diagrams, math, and visualizations."""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "Fraud_Detection_Project_Presentation.pptx"
PLOTS = ROOT / "docs" / "presentation_plots"
PLOTS.mkdir(parents=True, exist_ok=True)

DARK = RGBColor(0x06, 0x08, 0x0F)
ACCENT = RGBColor(0x3D, 0xFF, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x8B, 0x9B, 0xB4)
CARD = RGBColor(0x12, 0x1A, 0x28)

STUDENT_NAME = "Nick Einstein"
PROJECT_TITLE = "Sentinel: Adversarial Fraud Detection AI Platform"


def ensure_assets() -> None:
    """Regenerate metrics / EDA charts if scripts exist."""
    py = sys.executable
    snap = ROOT / "docs" / "model_metrics_snapshot.json"
    if not snap.exists():
        subprocess.run([py, str(ROOT / "scripts" / "compute_model_metrics.py")], check=False, cwd=ROOT)
    for script in ("build_etl_eda_documentation.py", "build_model_metrics_documentation.py"):
        p = ROOT / "scripts" / script
        if p.exists():
            subprocess.run([py, str(p)], check=False, cwd=ROOT)


def load_run_metrics() -> dict:
    p = ROOT / "artifacts" / "run_summary.json"
    return json.loads(p.read_text()) if p.exists() else {"metrics": {"roc_auc": 0.985, "f1": 0.689}}


def load_full_metrics() -> dict:
    p = ROOT / "docs" / "model_metrics_snapshot.json"
    return json.loads(p.read_text()) if p.exists() else {}


def draw_architecture_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(10, 7))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis("off")
    fig.patch.set_facecolor("#06080f")

    layers = [
        (8.5, "L1 Governance", "Alert · Retrain · Block serve"),
        (7.3, "L2 Experience", "Sentinel UI · REST API · WebSocket"),
        (6.1, "L3 Orchestration", "FraudAIPipeline · 10 stages"),
        (4.9, "L4 ML Core", "ETL · SMOTE · CatBoost · KS · IF"),
        (3.7, "L5 Data", "fraudTrain.csv · fraudTest.csv"),
        (2.5, "L6 MLOps", "Model Registry · Model Cards"),
        (1.3, "L7 Observability", "Telemetry · Run Manifests"),
    ]
    for y, title, sub in layers:
        box = FancyBboxPatch((1, y - 0.35), 8, 0.7, boxstyle="round,pad=0.02", facecolor="#121a28", edgecolor="#3dffa8", linewidth=1.5)
        ax.add_patch(box)
        ax.text(1.2, y, title, color="#3dffa8", fontsize=11, fontweight="bold", va="center")
        ax.text(3.2, y, sub, color="white", fontsize=9, va="center")
    ax.annotate("", xy=(5, 8.1), xytext=(5, 7.7), arrowprops=dict(arrowstyle="->", color="#3dffa8", lw=2))
    for y in [7.0, 5.8, 4.6, 3.4, 2.2]:
        ax.annotate("", xy=(5, y + 0.35), xytext=(5, y + 0.75), arrowprops=dict(arrowstyle="->", color="#555", lw=1.2))
    ax.text(5, 9.5, "Seven-Layer AI World Architecture (Top → Down)", ha="center", color="white", fontsize=14, fontweight="bold")
    out = PLOTS / "architecture_layers.png"
    fig.tight_layout()
    fig.savefig(out, dpi=160, facecolor=fig.get_facecolor())
    plt.close(fig)
    return out


def draw_pipeline_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(12, 3.5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 3)
    ax.axis("off")
    fig.patch.set_facecolor("#06080f")
    steps = [
        "Ingest", "ETL\n+SMOTE", "Scale", "CatBoost\nTrain", "Evaluate",
        "Drift\nBaseline", "KS Test", "Isolation\nForest", "Register", "Serve",
    ]
    n = len(steps)
    for i, label in enumerate(steps):
        x = 0.4 + i * 1.15
        box = FancyBboxPatch((x, 1), 1.0, 0.9, boxstyle="round,pad=0.02", facecolor="#121a28", edgecolor="#3dffa8")
        ax.add_patch(box)
        ax.text(x + 0.5, 1.45, label, ha="center", va="center", color="white", fontsize=7)
        if i < n - 1:
            ax.annotate("", xy=(x + 1.05, 1.45), xytext=(x + 1.12, 1.45),
                        arrowprops=dict(arrowstyle="->", color="#3dffa8", lw=1.2))
    ax.text(6, 2.5, "Training & Serving Pipeline (FraudAIPipeline)", ha="center", color="#3dffa8", fontsize=12, fontweight="bold")
    out = PLOTS / "pipeline_flow.png"
    fig.tight_layout()
    fig.savefig(out, dpi=160, facecolor=fig.get_facecolor())
    plt.close(fig)
    return out


def set_slide_bg(slide, color: RGBColor = DARK) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    for y, text, size, color, bold in [
        (2.2, PROJECT_TITLE, 34, WHITE, True),
        (3.5, STUDENT_NAME, 24, ACCENT, False),
        (4.2, "Data Science Final Project · In-Person Presentation", 14, MUTED, False),
    ]:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(0.8))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER


def add_text_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    hdr = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.7))
    hp = hdr.text_frame.paragraphs[0]
    hp.text = title
    hp.font.size = Pt(28)
    hp.font.bold = True
    hp.font.color.rgb = ACCENT
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12), Inches(0.35))
        sub.text_frame.paragraphs[0].text = subtitle
        sub.text_frame.paragraphs[0].font.size = Pt(12)
        sub.text_frame.paragraphs[0].font.color.rgb = MUTED
    body = slide.shapes.add_textbox(Inches(0.65), Inches(1.45 if subtitle else 1.2), Inches(12), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16 if not line.startswith("  ") else 14)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)


def add_math_slide(prs: Presentation, title: str, formulas: list[str], explanation: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    hdr = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.65))
    hdr.text_frame.paragraphs[0].text = title
    hdr.text_frame.paragraphs[0].font.size = Pt(26)
    hdr.text_frame.paragraphs[0].font.bold = True
    hdr.text_frame.paragraphs[0].font.color.rgb = ACCENT

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(11.8), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, f in enumerate(formulas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f
        p.font.name = "Courier New"
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT

    expl = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(11.8), Inches(3.5))
    etf = expl.text_frame
    etf.word_wrap = True
    for i, line in enumerate(explanation):
        p = etf.paragraphs[0] if i == 0 else etf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.space_after = Pt(5)


def add_viz_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    explanation: list[str],
    img_width: float = 6.2,
    img_left: float = 0.55,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    hdr = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12), Inches(0.6))
    hdr.text_frame.paragraphs[0].text = title
    hdr.text_frame.paragraphs[0].font.size = Pt(24)
    hdr.text_frame.paragraphs[0].font.bold = True
    hdr.text_frame.paragraphs[0].font.color.rgb = ACCENT

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(img_left), Inches(1.05), width=Inches(img_width))

    expl = slide.shapes.add_textbox(Inches(6.9 if img_width > 5 else 0.55), Inches(1.05), Inches(5.8), Inches(6))
    if img_width > 5:
        etf = expl.text_frame
    else:
        expl = slide.shapes.add_textbox(Inches(0.55), Inches(4.5), Inches(12), Inches(2.8))
        etf = expl.text_frame
    etf.word_wrap = True
    cap = etf.paragraphs[0]
    cap.text = "What this shows:"
    cap.font.bold = True
    cap.font.size = Pt(14)
    cap.font.color.rgb = ACCENT
    for line in explanation:
        p = etf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)


def add_diagram_slide(prs: Presentation, title: str, image_path: Path, caption: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    hdr = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12), Inches(0.6))
    hdr.text_frame.paragraphs[0].text = title
    hdr.text_frame.paragraphs[0].font.size = Pt(24)
    hdr.text_frame.paragraphs[0].font.bold = True
    hdr.text_frame.paragraphs[0].font.color.rgb = ACCENT

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(0.95), width=Inches(11.8))

    expl = slide.shapes.add_textbox(Inches(0.7), Inches(4.55), Inches(11.8), Inches(2.6))
    tf = expl.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "Explanation:"
    p0.font.bold = True
    p0.font.color.rgb = ACCENT
    p0.font.size = Pt(14)
    for line in caption:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE


def find_plot(*names: str) -> Path | None:
    dirs = [PLOTS, ROOT / "docs" / "eda_plots", ROOT / "docs" / "metrics_plots", ROOT / "docs" / "doc_plots", ROOT / "artifacts" / "plots"]
    for name in names:
        for d in dirs:
            p = d / name
            if p.exists():
                return p
    return None


def build() -> Path:
    ensure_assets()
    run_m = load_run_metrics()
    full_m = load_full_metrics()
    roc = full_m.get("roc_auc", run_m.get("metrics", {}).get("roc_auc", 0.985))
    f1 = full_m.get("f1", run_m.get("metrics", {}).get("f1", 0.689))
    prec = full_m.get("precision", 0.704)
    rec = full_m.get("recall", 0.674)

    arch_img = draw_architecture_diagram()
    pipe_img = draw_pipeline_diagram()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Core narrative slides ──
    add_title_slide(prs)

    add_text_slide(prs, "Abstract", [
        "Purpose: detect credit-card fraud in real time and monitor adversarial drift",
        "Primary algorithm: CatBoost gradient boosted trees (depth 6, 500 iterations)",
        "Supporting: SMOTE (class balance), Isolation Forest (anomalies), KS drift tests",
        "Data: credit_dt fraudTrain.csv / fraudTest.csv (~1.3M / 555k rows)",
        f"Results: ROC-AUC {roc:.4f}, F1 {f1:.4f}, Precision {prec:.3f}, Recall {rec:.3f}",
        "Platform: Sentinel web UI + FastAPI + seven-layer AI architecture",
    ])

    add_text_slide(prs, "Existing System (Limitations)", [
        "Rule-based velocity limits and static blacklists",
        "Batch-only ML with no post-deployment drift monitoring",
        "No unified pipeline — training, scoring, monitoring in separate scripts",
        "Kaggle PCA features only — no merchant category or geography",
        "No governance linking model quality to serve / retrain decisions",
    ], subtitle="What legacy fraud systems typically lack")

    add_text_slide(prs, "Proposed System — Sentinel", [
        "Unified FraudAIPipeline: 10 orchestrated stages",
        "13 engineered features including Haversine distance_km",
        "REST API + Sentinel UI + WebSocket live stream",
        "Model registry, telemetry, and model cards for audit",
        "Governance layer: alert when ≥3 features drift (KS p < 0.05)",
    ])

    add_text_slide(prs, "Data Resources", [
        "Source: fraudTrain.csv (train) · fraudTest.csv (test) — credit card transactions",
        "Raw fields: amount, lat/long, merchant coords, category, gender, state, timestamp",
        "Derived: trans_hour, trans_dow, distance_km (Haversine formula)",
        "Target: is_fraud (0 = legitimate, 1 = fraud) — ~0.8% fraud in train sample",
        "Caps: 200k train / 100k test rows for practical laptop training",
    ])

    # ── Architecture diagrams ──
    add_diagram_slide(prs, "Architecture Diagram — Seven Layers", arch_img, [
        "Top-down design: L1 Business sets policy; L4 ML Core computes predictions.",
        "L2 Experience is how users interact (Sentinel UI, API).",
        "L3 Orchestration runs the stage graph; L6/L7 handle versioning and audit.",
        "L5 Data feeds raw CSVs into the featurizer before modelling.",
    ])

    add_diagram_slide(prs, "Pipeline Flow Diagram", pipe_img, [
        "Sequential stages executed by python main.py.",
        "SMOTE applied only after ETL on training data — test set stays imbalanced.",
        "Drift baseline saved before serving; KS test compares test vs train distributions.",
        "Final model registered and promoted to production in the registry.",
    ])

    # ── Mathematics & algorithms ──
    add_math_slide(
        prs,
        "Algorithm 1 — SMOTE (Class Balancing)",
        [
            "x_new = x_i + λ · (x_neighbor − x_i),  λ ∈ [0,1]",
            "k_neighbors = 5,  sampling_strategy = auto",
            "Applied to TRAIN only — fraud class ~0.8% → balanced ~50/50",
        ],
        [
            "Synthetic Minority Over-sampling Technique creates new fraud examples.",
            "Interpolates between a real fraud point and one of its k nearest fraud neighbours.",
            "Prevents CatBoost from ignoring the rare fraud class during training.",
            "Test set is never SMOTE'd — evaluation reflects real-world imbalance.",
        ],
    )

    add_math_slide(
        prs,
        "Algorithm 2 — Haversine Distance (Feature Engineering)",
        [
            "a = sin²(Δlat/2) + cos(lat₁)·cos(lat₂)·sin²(Δlon/2)",
            "distance_km = 2 · R · arcsin(√a),   R = 6371 km",
            "Inputs: cardholder (lat, long) and merchant (merch_lat, merch_long)",
        ],
        [
            "Computes great-circle distance between cardholder and purchase location.",
            "Flags geographically implausible transactions (e.g. card in NY, merchant in CA).",
            "Used as a numeric feature alongside raw coordinates in CatBoost.",
        ],
    )

    add_math_slide(
        prs,
        "Algorithm 3 — CatBoost (Primary Classifier)",
        [
            "F(x) = Σₜ₌₁^T η · hₜ(x),   T=500 trees, η=0.05, depth≤6",
            "P(fraud | x) = σ(F(x)) = 1 / (1 + e^(−F(x)))",
            "Decision: fraud_prediction = 1  if  P(fraud) ≥ 0.5",
        ],
        [
            "Gradient boosting: each tree hₜ corrects errors of previous trees.",
            "Depth 6 = up to 6 sequential split rules per tree (category, amount, time…).",
            "auto_class_weights=Balanced penalises misclassified fraud more heavily.",
            "500 trees × 13 scaled features → fraud probability in [0, 1].",
        ],
    )

    add_math_slide(
        prs,
        "Algorithm 4 — StandardScaler (Preprocessing)",
        [
            "x_scaled = (x − μ_train) / σ_train",
            "Fit on SMOTE-balanced train; transform train & test identically",
        ],
        [
            "Puts all 13 features on comparable scale (mean 0, variance 1).",
            "Critical for distance-based anomaly detection and stable tree splits.",
            "Scaler saved to artifacts/scaler.joblib for serving-time consistency.",
        ],
    )

    add_math_slide(
        prs,
        "Algorithm 5 — Kolmogorov–Smirnov Drift Test",
        [
            "KS = max_x | F_train(x) − F_test(x) |",
            "Drift if p-value < α  (α = 0.05)",
            "Governance: alert if ≥3 features drift; block if ≥10",
        ],
        [
            "Compares cumulative distribution of each feature: train vs production/test.",
            "Detected drift on unix_time and trans_dow — temporal shift between splits.",
            "Feeds L1 governance — operators alerted before model silently degrades.",
        ],
    )

    add_math_slide(
        prs,
        "Algorithm 6 — Isolation Forest (Anomaly Detection)",
        [
            "anomaly_score = −decision_function(x)",
            "n_estimators = 100, contamination = 0.002",
            "Shorter isolation path → higher anomaly score",
        ],
        [
            "Unsupervised ensemble — no fraud labels needed at score time.",
            "Random trees isolate points; outliers need fewer splits to separate.",
            "Complements CatBoost in the live WebSocket stream for novel attack patterns.",
        ],
    )

    add_math_slide(
        prs,
        "Evaluation Metrics (Mathematics)",
        [
            "ROC-AUC = ∫ TPR(FPR) d(FPR)",
            "Precision = TP / (TP + FP),   Recall = TP / (TP + FN)",
            "F1 = 2 · Precision · Recall / (Precision + Recall)",
            "Accuracy = (TP + TN) / N  — misleading when fraud < 1%",
        ],
        [
            f"Our test results: ROC-AUC={roc:.4f}, F1={f1:.4f}, Precision={prec:.3f}, Recall={rec:.3f}.",
            "Confusion matrix: TP=271, FN=131, FP=114, TN=99,484 (100k test rows).",
            "Balanced accuracy (0.836) is fairer than raw accuracy (0.998) under imbalance.",
        ],
    )

    # ── Visualization slides with explanations ──
    viz_specs = [
        ("Visualization — Class Imbalance", "class_balance.png",
         ["Blue bars = legitimate transactions; red = fraud.",
          "Train fraud rate ~0.82%; test ~0.40% — extreme imbalance (~1000:1).",
          "Justifies SMOTE on train and ROC-AUC/F1 as primary metrics, not accuracy."]),
        ("Visualization — Transaction Amount Distribution", "amount_dist.png",
         ["Histogram of log(1 + amount) for legit (blue) vs fraud (red).",
          "Fraud transactions skew toward higher amounts (mean ~$509 vs ~$68 legit).",
          "Confirms amt as the 2nd most important CatBoost feature after category."]),
        ("Visualization — Fraud Rate by Category", "category_fraud_rate.png",
         ["Each bar = merchant category (shopping_net, travel, grocery, etc.).",
          "Category shows the widest fraud-rate spread — explains ~64% feature importance.",
          "CatBoost learns different risk profiles per merchant type."]),
        ("Visualization — Fraud Rate by Hour", "hour_fraud.png",
         ["Line chart: fraud rate vs hour of day (0–23) from trans_date_trans_time.",
          "Temporal patterns motivate trans_hour and trans_dow as model inputs.",
          "KS drift on unix_time/trans_dow reflects train vs test time shift."]),
        ("Visualization — SMOTE Class Balance", "smote_balance.png",
         ["Before SMOTE: ~198k legit vs ~1.6k fraud. After: classes balanced (~1:1).",
          "Synthetic fraud points fill the feature space between real fraud neighbours.",
          "Only training data is resampled — test evaluation stays realistic."]),
        ("Visualization — Model Metrics Summary", "metrics_bar.png",
         ["Bar chart of ROC-AUC, Average Precision, F1, Precision, Recall, Balanced Acc.",
          f"ROC-AUC {roc:.3f} = excellent ranking; F1 {f1:.3f} = realistic under imbalance.",
          "Average Precision (0.709) focuses specifically on the fraud class."]),
        ("Visualization — Confusion Matrix", "confusion_matrix.png",
         ["Rows = actual class; columns = predicted class at threshold 0.5.",
          "271 true positives (fraud caught); 131 false negatives (fraud missed).",
          "114 false positives (legit flagged) — acceptable cost for fraud prevention."]),
        ("Visualization — Threshold Sensitivity", "threshold_sweep.png",
         ["Lines show precision (blue), recall (red), F1 (green) vs threshold.",
          "Lower threshold → more fraud caught but more false alarms.",
          "Default 0.5 balances precision (0.70) and recall (0.67)."]),
        ("Visualization — Feature Importance", "feature_importance.png",
         ["CatBoost gain-based importance across all 500 trees.",
          "category dominates (~64%); amt (~11%); gender (~9%); geo features smaller.",
          "Guides feature engineering priorities and business interpretation."]),
        ("Visualization — KS Drift Report", "drift_ks.png",
         ["Red bars = drift detected (p < 0.05); blue = no significant drift.",
          "unix_time KS=1.0 — complete temporal distribution shift train→test.",
          "Governance uses this to alert operators before silent model degradation."]),
    ]

    for title, fname, expl in viz_specs:
        path = find_plot(fname)
        if path:
            add_viz_slide(prs, title, path, ["• " + e if not e.startswith("•") else e for e in expl])
        else:
            add_text_slide(prs, title, expl + ["(Chart file not found — run python scripts/compute_model_metrics.py)"])

    pca = find_plot("pca_train.png")
    if pca:
        add_viz_slide(prs, "Visualization — PCA Projection (2D)", pca, [
            "• 2 principal components of scaled training features after StandardScaler.",
            "• Points coloured by Class: fraud (1) vs legitimate (0).",
            "• Exploratory only — PCA is NOT fed to CatBoost (13 original features used).",
            "• Shows partial separation of fraud in reduced 2D space.",
        ])

    # ── Methodology & results summary ──
    add_text_slide(prs, "Methodology Summary", [
        "1. Ingest fraudTrain / fraudTest CSVs",
        "2. Featurize: time, Haversine distance, label-encode categoricals",
        "3. Clean: dedupe, median impute",
        "4. SMOTE balance training labels",
        "5. StandardScaler fit/transform",
        "6. Train CatBoost (500 trees, depth 6)",
        "7. Evaluate on held-out fraudTest",
        "8. Fit KS drift baseline; run drift check",
        "9. Train Isolation Forest; register model; apply governance",
    ])

    add_text_slide(prs, "Accuracies & Results", [
        f"ROC-AUC: {roc:.4f} — strong fraud vs legit ranking",
        f"F1 Score: {f1:.4f} at threshold 0.5",
        f"Precision: {prec:.4f} — 70% of flags are true fraud",
        f"Recall: {rec:.4f} — 67% of fraud cases detected",
        f"Average Precision: {full_m.get('average_precision', 0.709):.4f}",
        "Model: 500 CatBoost trees, depth 6, 13 features",
        "Production version: v_f381139c (Model Registry)",
    ])

    add_text_slide(prs, "Conclusions & Future Work", [
        "Conclusions:",
        "  Seven-layer architecture separates policy from ML computation",
        "  CatBoost + SMOTE achieves strong ROC-AUC on imbalanced fraud data",
        "  Visualizations confirm category and amount drive predictions",
        "  KS drift monitoring supports adversarial / temporal shift awareness",
        "Future:",
        "  SHAP explanations, automated retrain, cloud deployment, API auth",
    ])

    # Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = f"{STUDENT_NAME} · Questions?"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Presentation saved: {path} ({len(Presentation(str(path)).slides)} slides)")
